import io
import PyPDF2
from docx import Document
from pptx import Presentation
from fastapi import HTTPException

def extract_text_from_pdf(file_content: bytes) -> str:
    """
    Extract text from PDF file
    
    Args:
        file_content: PDF file content as bytes
    
    Returns:
        Extracted text as string
    """
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text = ""
        
        for page in pdf_reader.pages:
            try:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            except Exception as e:
                print(f"Error extracting text from PDF page: {e}")
                continue
        
        if not text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the PDF")
        
        return text.strip()
        
    except PyPDF2.errors.PdfReadError as e:
        raise HTTPException(status_code=400, detail=f"Invalid PDF file: {str(e)}")
    except Exception as e:
        print(f"Error processing PDF: {e}")
        raise HTTPException(status_code=500, detail="Failed to process PDF file")

def extract_text_from_docx(file_content: bytes) -> str:
    """
    Extract text from DOCX file
    
    Args:
        file_content: DOCX file content as bytes
    
    Returns:
        Extracted text as string
    """
    try:
        doc = Document(io.BytesIO(file_content))
        text = ""
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + "\n"
        
        if not text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the DOCX file")
        
        return text.strip()
        
    except Exception as e:
        print(f"Error processing DOCX: {e}")
        raise HTTPException(status_code=500, detail="Failed to process DOCX file")

def extract_text_from_pptx(file_content: bytes) -> str:
    """
    Extract text from PPTX file
    
    Args:
        file_content: PPTX file content as bytes
    
    Returns:
        Extracted text as string
    """
    try:
        presentation = Presentation(io.BytesIO(file_content))
        text = ""
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = f"Slide {slide_num}:\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
                
                # Extract text from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_text += cell.text + "\n"
            
            if slide_text.strip() != f"Slide {slide_num}:":
                text += slide_text + "\n"
        
        if not text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the PPTX file")
        
        return text.strip()
        
    except Exception as e:
        print(f"Error processing PPTX: {e}")
        raise HTTPException(status_code=500, detail="Failed to process PPTX file")

def extract_text_from_file(file_content: bytes, filename: str) -> str:
    """
    Extract text from file based on file extension
    
    Args:
        file_content: File content as bytes
        filename: Name of the file (used to determine file type)
    
    Returns:
        Extracted text as string
    """
    file_extension = filename.lower().split('.')[-1]
    
    if file_extension == 'pdf':
        return extract_text_from_pdf(file_content)
    elif file_extension == 'docx':
        return extract_text_from_docx(file_content)
    elif file_extension == 'pptx':
        return extract_text_from_pptx(file_content)
    else:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported file type: .{file_extension}. Supported types: PDF, DOCX, PPTX"
        )

def validate_file_type(filename: str) -> bool:
    """
    Validate if file type is supported
    
    Args:
        filename: Name of the file
    
    Returns:
        True if file type is supported, False otherwise
    """
    allowed_extensions = ['pdf', 'docx', 'pptx']
    file_extension = filename.lower().split('.')[-1]
    return file_extension in allowed_extensions

def validate_file_size(file_size: int) -> bool:
    """
    Validate file size (10KB - 10MB)
    
    Args:
        file_size: Size of the file in bytes
    
    Returns:
        True if file size is valid, False otherwise
    """
    min_size = 10_000  # 10KB
    max_size = 10_000_000  # 10MB
    return min_size <= file_size <= max_size